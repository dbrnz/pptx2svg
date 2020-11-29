#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import colorsys
import json
import os
import re
import string

from collections import OrderedDict
from math import sqrt, sin, cos, pi as PI
from zipfile import ZipFile

#===============================================================================

import numpy as np
import svgwrite
from tqdm import tqdm

#===============================================================================

import pptx.shapes.connector
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE, MSO_THEME_COLOR, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Length

#===============================================================================

from formula import Geometry, radians
from presets import DML, ThemeDefinition

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

POINTS_PER_IN = 72

# SVG pixel resolution
PIXELS_PER_IN = 96
EMU_PER_PIXEL = EMU_PER_IN/PIXELS_PER_IN

# Minimum width for a stroked path in points
MIN_STROKE_WIDTH = 0.5

#===============================================================================

def emu_to_pixels(emu):
#======================
    return emu/EMU_PER_PIXEL

def points_to_pixels(pts):
#=========================
    return pts*PIXELS_PER_IN/POINTS_PER_IN

def transform_point(transform, point):
#=====================================
    return (transform@[point[0], point[1], 1.0])[:2]

#===============================================================================

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

ARROW_MARKERS = {
    'triangle-head': 'M 10 0 L 0 5 L 10 10 z',
    'triangle-tail': 'M 0 0 L 10 5 L 0 10 z'
}

## NB. Adobe Illustrator 2020 doesn't appear to support marker definitions in SVG

def add_marker_definitions(drawing):
#===================================
    # arrowhead markers (see https://developer.mozilla.org/en-US/docs/Web/SVG/Element/marker)
    for id, path in ARROW_MARKERS.items():
        marker = drawing.marker(id=id,
                                viewBox="0 0 10 10",
                                refX="5", refY="5",
                                markerWidth="6", markerHeight="6",
                                orient="auto")
        marker.add(drawing.path(d=path))
        drawing.defs.add(marker)

def marker_id(marker_def, end):
#==============================
    marker_type = marker_def.get('type')
    return ('#{}-{}'.format(marker_type, end)
            if marker_type is not None
            else None)

#===============================================================================

# Don't set a path id for default shape names

EXCLUDED_NAME_PREFIXES = [
    'Freeform',
    'Group',
    'Oval',
]

# Markup that has been deprecated

EXCLUDED_NAME_MARKUP = [
    '.siblings',
]

# Check to see if we have a valid name and encode it as an id

def id_from_name(name):
#======================
    if name not in EXCLUDED_NAME_MARKUP:
        for prefix in EXCLUDED_NAME_PREFIXES:
            if name.startswith(prefix):
                return None
        return adobe_encode(name)
    return None

# Helpers for encoding names for Adobe Illustrator

def match_to_hex(m):
#===================
    c = m[0]
    return (c   if c in (string.ascii_letters + string.digits) else
            '_' if c in string.whitespace else
            '_x{:02X}_'.format(ord(c)))

def adobe_encode(s):
#===================
    return re.sub('.', match_to_hex, s)

#===============================================================================

class Theme(object):
    def __init__(self, pptx_source):
        with ZipFile(pptx_source, 'r') as presentation:
            for info in presentation.infolist():
                if info.filename.startswith('ppt/theme/'):
                    self.__theme_definition = ThemeDefinition.new(presentation.read(info))
                    break

    def colour_scheme(self):
    #=======================
        return self.__theme_definition.themeElements.clrScheme

#===============================================================================

class ColourMap(object):
    def __init__(self, ppt_theme, slide):
        self.__colour_defs = {}
        for colour_def in ppt_theme.colour_scheme():
            defn = colour_def[0]
            if defn.tag == DML('sysClr'):
                self.__colour_defs[colour_def.tag] = RGBColor.from_string(defn.attrib['lastClr'])
            elif defn.tag == DML('srgbClr'):
                self.__colour_defs[colour_def.tag] = RGBColor.from_string(defn.val)
        # The slide's layout master can have colour aliases
        colour_map = slide.slide_layout.slide_master.element.clrMap.attrib
        for key, value in colour_map.items():
            if key != value:
                self.__colour_defs[DML(key)] = self.__colour_defs[DML(value)]

    def lookup(self, colour_format):
    #===============================
        if colour_format.type == MSO_COLOR_TYPE.RGB:
            rgb = colour_format.rgb
        elif colour_format.type == MSO_COLOR_TYPE.SCHEME:
            key = MSO_THEME_COLOR.to_xml(colour_format.theme_color)
            rgb = self.__colour_defs[DML(key)]
        else:
            raise ValueError('Unsupported colour format: {}'.format(colour_format.type))
        brightness = colour_format.brightness
        if brightness != 0.0:
            hsv = list(colorsys.rgb_to_hsv(*(np.array(rgb)/255.0)))
            hsv[2] *= (brightness + 1.0)
            if hsv[2] > 1.0:
                hsv[2] = 1.0
            colour = np.uint8(255*np.array(colorsys.hsv_to_rgb(*hsv)) + 0.5)
            rgb = RGBColor(*colour.tolist())
        return '#{}'.format(str(rgb))

#===============================================================================

class Gradient(object):
    def __init__(self, id, fill, colour_map):
        self.__id = id
        self.__kind = 'Linear'
        self.__angle = fill.gradient_angle
        self.__stops = [(stop.position, colour_map.lookup(stop.color))
                            for stop in fill.gradient_stops]

    @property
    def id(self):
        return self.__id

    def svg_definition(self):
    #========================
        stops = ['<stop offset="{}%" stop-color="{}"/>'.format(100*stop[0], stop[1])
                    for stop in self.__stops]
        return ('<{gradient}Gradient id="{id}">{stops}</{gradient}Gradient>'
                .format(gradient=self.__kind, id=self.__id, stops='/n'.join(stops)))

## WIP  Want list of unique gradient definitions

#===============================================================================

class DrawMLTransform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = ((xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = ((xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    bbox)
        (Bx_, By_) = (xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = (xfrm.ext.cx, xfrm.ext.cy)
        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1
        T_st = np.array([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                         [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                         [     0,      0,                 1]])
        U = np.array([[1, 0, -(Bx_ + Dx_/2.0)],
                      [0, 1, -(By_ + Dy_/2.0)],
                      [0, 0,                1]])
        R = np.array([[cos(theta), -sin(theta), 0],
                      [sin(theta),  cos(theta), 0],
                      [0,                    0, 1]])
        Flip = np.array([[Fx,  0, 0],
                         [ 0, Fy, 0],
                         [ 0,  0, 1]])
        T_rf = np.linalg.inv(U)@R@Flip@U
        self.__T = T_rf@T_st

    def matrix(self):
        return self.__T

#===============================================================================

class SvgLayer(object):
    def __init__(self, size, slide, slide_number, ppt_theme, quiet=False):
        self.__slide = slide
        self.__colour_map = ColourMap(ppt_theme, slide)
        self.__dwg = svgwrite.Drawing(filename=None, size=size)
## WIP  add_marker_definitions(self.__dwg)
        self.__id = None
        self.__models = None
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                for part in notes_text[1:].split():
                    id_match = re.match('id *\((.*)\)', part)
                    if id_match is not None:
                        self.__id = id_match[1].strip()
                    models_match = re.match('models *\((.*)\)', part)
                    if models_match is not None:
                        self.__models = models_match[1].strip()
        if self.__id is None:
            self.__id = 'slide-{:02d}'.format(slide_number)
        self.__filename = None
        self.__quiet =  quiet

    @property
    def filename(self):
        return self.__filename

    @property
    def id(self):
        return self.__id

    @property
    def models(self):
        return self.__models

    def save(self, output_dir):
    #==========================
        self.__filename = os.path.join(output_dir, '{}.svg'.format(self.__id))
        self.__dwg.saveas(self.__filename, pretty=True, indent=4)

    def process(self, transform):
    #============================
        self.process_shape_list(self.__slide.shapes, self.__dwg, transform,  not self.__quiet)

    def process_group(self, group, svg_parent, transform):
    #=====================================================
        svg_group = self.__dwg.g(id=id_from_name(group.name))
        svg_parent.add(svg_group)
        self.process_shape_list(group.shapes, svg_group, transform@DrawMLTransform(group).matrix())

    def process_shape_list(self, shapes, svg_parent, transform, show_progress=False):
    #================================================================================
        if show_progress:
            print('Processing shape list...')
            progress_bar = tqdm(total=len(shapes),
                unit='shp', ncols=40,
                bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                self.process_shape(shape, svg_parent, transform)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.process_group(shape, svg_parent, transform)
            elif (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
               or shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            if show_progress:
                progress_bar.update(1)
        if show_progress:
            progress_bar.close()

    def process_shape(self, shape, svg_parent, transform):
    #=====================================================
        id = id_from_name(shape.name)
        geometry = Geometry(shape)
        if id is not None and len(geometry) > 1:
            # Add a group to hold multiple paths
            ## We should really add a `.group` placeholder
            group = self.__dwg.g(id=id)
            svg_parent.add(group)
            svg_parent = group
            id = None

        for path in geometry.path_list:
            svg_path = self.__dwg.path(fill='none', class_='non-scaling-stroke')
            if id is not None:
                svg_path.attribs['id'] = id
            bbox = (shape.width, shape.height) if path.w is None else (path.w, path.h)
            T = transform@DrawMLTransform(shape, bbox).matrix()
            first_point = None
            current_point = None
            closed = False
            for c in path.getchildren():
                if   c.tag == DML('arcTo'):
                    wR = geometry.attrib_value(c, 'wR')
                    hR = geometry.attrib_value(c, 'hR')
                    stAng = radians(geometry.attrib_value(c, 'stAng'))
                    swAng = radians(geometry.attrib_value(c, 'swAng'))
                    p1 = ellipse_point(wR, hR, stAng)
                    p2 = ellipse_point(wR, hR, stAng + swAng)
                    pt = (current_point[0] - p1[0] + p2[0],
                          current_point[1] - p1[1] + p2[1])
                    large_arc_flag = 1 if swAng >= PI else 0
                    svg_path.push('A', emu_to_pixels(wR), emu_to_pixels(hR),
                                       0, large_arc_flag, 1,
                                       *transform_point(T, pt))
                    current_point = pt

                elif c.tag == DML('close'):
                    if first_point is not None and current_point != first_point:
                        svg_path.push('Z')
                    closed = True
                    first_point = None
                elif c.tag == DML('cubicBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.extend(transform_point(T, pt))
                        current_point = pt
                    svg_path.push('C', *coords)
                elif c.tag == DML('lnTo'):
                    pt = geometry.point(c.pt)
                    coords = transform_point(T, pt)
                    svg_path.push('L', *coords)
                    current_point = pt
                elif c.tag == DML('moveTo'):
                    pt = geometry.point(c.pt)
                    coords = transform_point(T, pt)
                    svg_path.push('M', *coords)
                    if first_point is None:
                        first_point = pt
                    current_point = pt
                elif c.tag == DML('quadBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.extend(transform_point(T, pt))
                        current_point = pt
                    svg_path.push('Q', *coords)
                else:
                    print('Unknown path element: {}'.format(c.tag))

            if closed:
                if shape.fill.type == MSO_FILL_TYPE.SOLID:
                    svg_path.attribs['fill'] = self.__colour_map.lookup(shape.fill.fore_color)
                    alpha = shape.fill.fore_color.alpha
                    if alpha < 1.0:
                        svg_path.attribs['opacity'] = alpha
                elif shape.fill.type == MSO_FILL_TYPE.GRADIENT:
## WIP              gradient = Gradient('id', shape.fill, self.__colour_map)
                    svg_path.attribs['fill'] = '#CCCCCC'   ## TEMP
                    svg_path.attribs['opacity'] = 0.3      ## TEMP
                elif shape.fill.type is None:
                    svg_path.attribs['fill'] = '#000000'
                elif shape.fill.type != MSO_FILL_TYPE.BACKGROUND:
                    print('Unsupported fill type: {}'.format(shape.fill.type))

            if shape.line.fill.type == MSO_FILL_TYPE.SOLID:
                svg_path.attribs['stroke'] = self.__colour_map.lookup(shape.line.color)
                alpha = shape.line.fill.fore_color.alpha
                if alpha < 1.0:
                    svg_path.attribs['stroke-opacity'] = alpha
            elif shape.line.fill.type is None:
                svg_path.attribs['stroke'] = '#000000'
            elif shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:
                print('Unsupported line fill type: {}'.format(shape.line.fill.type))

            stroke_width = points_to_pixels(max(Length(shape.line.width).pt, MIN_STROKE_WIDTH))
            svg_path.attribs['stroke-width'] = stroke_width
            if shape.line.dash_style is not None:
                if shape.line.dash_style == MSO_LINE_DASH_STYLE.DASH:
                    svg_path.attribs['stroke-dasharray'] = 4*stroke_width
                elif shape.line.dash_style == MSO_LINE_DASH_STYLE.DASH_DOT:
                    svg_path.attribs['stroke-dasharray'] = '{} {} {} {}'.format(4*stroke_width, stroke_width, stroke_width, stroke_width)
                elif shape.line.dash_style == MSO_LINE_DASH_STYLE.LONG_DASH:
                    svg_path.attribs['stroke-dasharray'] = '{} {}'.format(4*stroke_width, stroke_width)
                elif shape.line.dash_style == MSO_LINE_DASH_STYLE.SQUARE_DOT:
                    svg_path.attribs['stroke-dasharray'] = '{} {}'.format(2*stroke_width, stroke_width)
                elif shape.line.dash_style != MSO_LINE_DASH_STYLE.SOLID:
                    print('Unsupported line dash style: {}'.format(shape.line.dash_style))

## WIP      if 'type' in shape.line.headEnd or 'type' in shape.line.tailEnd:
## WIP          svg_path.set_markers((marker_id(shape.line.headEnd, 'head'),
## WIP                                None,
## WIP                                marker_id(shape.line.tailEnd, 'tail')))

            svg_parent.add(svg_path)

#===============================================================================

class SvgExtractor(object):
    def __init__(self, options):
        self.__pptx = Presentation(options.powerpoint)
        self.__theme = Theme(options.powerpoint)
        self.__slides = self.__pptx.slides
        (pptx_width, pptx_height) = (self.__pptx.slide_width, self.__pptx.slide_height)
        self.__transform = np.array([[1.0/EMU_PER_PIXEL,                 0, 0],
                                     [                0, 1.0/EMU_PER_PIXEL, 0],
                                     [                0,                 0, 1]])
        self.__svg_size = transform_point(self.__transform, (pptx_width, pptx_height))
        self.__output_dir = options.output_dir
        self.__debug = options.debug
        self.__quiet = options.quiet
        self.__saved_svg = OrderedDict()
        self.__id = None
        self.__models = None

    @property
    def id(self):
        return self.__id

    @property
    def saved_svg(self):
        return self.__saved_svg

    def slide_to_svg(self, slide, slide_number):
    #===========================================
        if self.__debug:
            with open(os.path.join(self.__output_dir, 'slide-{:02d}.xml'.format(slide_number)), 'w') as xml:
                xml.write(slide.element.xml)
        layer = SvgLayer(self.__svg_size, slide, slide_number, self.__theme, self.__quiet)
        layer.process(self.__transform)
        layer.save(self.__output_dir)
        self.__saved_svg[layer.id] = layer.filename
        if slide_number == 1:
            if not layer.id.startswith('slide-'):
                self.__id = layer.id
            self.__models = layer.models

    def slides_to_svg(self):
    #=======================
        for n, slide in enumerate(self.__slides):
            self.slide_to_svg(slide, n+1)

    def manifest(self):
    #==================
        manifest = OrderedDict()
        if self.__id is not None:
            manifest['id'] = self.__id
        if self.__models is not None:
            manifest['models'] = self.__models
        manifest['sources'] = []
        source_kind = 'base'
        for id, filename in self.__saved_svg.items():
            manifest['sources'].append(OrderedDict(
                id=id,
                href=filename,
                kind=source_kind
            ))
            if source_kind == 'base':
                source_kind = 'details'
        return manifest

    def save_manifest(self, filename):
    #=================================
        with open(filename, 'w') as output:
            output.write(json.dumps(extractor.manifest(), indent=4))

#===============================================================================

if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Convert Powerpoint slides to SVG.')

    parser.add_argument('-d', '--debug', action='store_true', help='save DrawML to aid with debugging')
    parser.add_argument('-q', '--quiet', action='store_true', help='do not show progress bar')

    parser.add_argument('powerpoint', metavar='POWERPOINT_FILE',
                        help='the Powerpoint file to convert')

    parser.add_argument('output_dir', metavar='OUTPUT_DIRECTORY',
                        help='directory in which to save SVG files')

    args = parser.parse_args()

    if not os.path.exists(args.output_dir):
        os.makedirs(args.output_dir)

    extractor = SvgExtractor(args)
    extractor.slides_to_svg()

    manifest = os.path.join(args.output_dir, 'manifest.json')
    extractor.save_manifest(manifest)
    print('Manifest saved as `{}`'.format(manifest))

#===============================================================================

